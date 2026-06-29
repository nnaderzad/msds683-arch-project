"""Build the MSDS683 final-presentation deck as a .pptx (imports cleanly into Google Slides).

Deterministic + re-runnable. Content mirrors final_presentation/PLAN.md. This is a ROUGH
DRAFT skeleton: titles, bullets, a pipeline diagram, and detailed speaker notes (timing +
rubric reminders) live in the notes pane of each slide.

Run (dedicated env, NOT a project env):
    conda activate presentation        # see final_presentation/requirements.txt
    python final_presentation/build_deck.py

Output: final_presentation/event-demand-final-deck.pptx
Import: upload the .pptx to Google Drive and open with Google Slides (or Slides > File >
Import slides). Speaker notes carry over.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- palette / layout
NAVY = RGBColor(0x1F, 0x2D, 0x4E)
ACCENT = RGBColor(0x2E, 0x86, 0xAB)
INK = RGBColor(0x22, 0x26, 0x2B)
MUTED = RGBColor(0x5B, 0x63, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)

LAYER = {  # medallion colour-coding for the pipeline diagram
    "sources": RGBColor(0x4C, 0x72, 0xB0),
    "bronze": RGBColor(0xA9, 0x6B, 0x3F),
    "silver": RGBColor(0x8A, 0x90, 0x99),
    "gold": RGBColor(0xD2, 0xA3, 0x2E),
    "app": RGBColor(0x3F, 0x9B, 0x5E),
}

SW, SH = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = slide.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    return sp


def _text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs: list of paragraphs, each a list of (text, size, color, bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(6)
        for text, size, color, bold in para:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Arial"
    return tb


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()


def _content_slide(prs, kicker, title, bullets, notes):
    """Standard content slide: accent kicker + title bar, bulleted body, speaker notes.

    bullets: list of (text, level) where level 0 = main, 1 = sub.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    # title band
    _box(slide, 0, 0, SW, Inches(1.35), fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    _box(slide, 0, Inches(1.35), SW, Inches(0.06), fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, MARGIN, Inches(0.18), SW - 2 * MARGIN, Inches(0.4),
          [[(kicker.upper(), 13, RGBColor(0x9F, 0xC4, 0xD8), True)]])
    _text(slide, MARGIN, Inches(0.5), SW - 2 * MARGIN, Inches(0.8),
          [[(title, 28, WHITE, True)]])
    # body bullets
    paras = []
    for text, level in bullets:
        glyph = "•  " if level == 0 else "      –  "
        size = 18 if level == 0 else 15
        color = INK if level == 0 else MUTED
        paras.append([(glyph + text, size, color, False)])
    _text(slide, MARGIN, Inches(1.7), SW - 2 * MARGIN, Inches(5.4), paras)
    _notes(slide, notes)
    return slide


# ---------------------------------------------------------------- individual slides
def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, NAVY)
    _box(slide, 0, Inches(3.55), SW, Inches(0.06), fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, MARGIN, Inches(2.2), SW - 2 * MARGIN, Inches(1.4),
          [[("Event-Demand Forecasting", 44, WHITE, True)]])
    _text(slide, MARGIN, Inches(3.75), SW - 2 * MARGIN, Inches(0.8),
          [[("A medallion data architecture for live-music demand signals", 22,
             RGBColor(0xC9, 0xD6, 0xE5), False)]])
    _text(slide, MARGIN, Inches(6.4), SW - 2 * MARGIN, Inches(0.6),
          [[("MSDS683 · Final Presentation & Demo · [team names]", 16,
             RGBColor(0x9F, 0xC4, 0xD8), False)]])
    _notes(slide, """
~10 min total. Open fast — assume the audience saw the midterm. One line: "We built a
bronze->silver->gold pipeline on GCP that turns three live data feeds into one model-ready
table and a consumable app." Then move straight to the recap. [Fill in team names + who
presents which section.]
""")


def pipeline_slide(prs):
    """Slide 4 — architecture map: Sources -> Bronze -> Silver -> Gold -> App."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, WHITE)
    _box(slide, 0, 0, SW, Inches(1.35), fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    _box(slide, 0, Inches(1.35), SW, Inches(0.06), fill=ACCENT, shape=MSO_SHAPE.RECTANGLE)
    _text(slide, MARGIN, Inches(0.18), SW - 2 * MARGIN, Inches(0.4),
          [[("PIPELINE -> END PRODUCT  ·  ~5 MIN  ·  the core (25 pts)", 13,
             RGBColor(0x9F, 0xC4, 0xD8), True)]])
    _text(slide, MARGIN, Inches(0.5), SW - 2 * MARGIN, Inches(0.8),
          [[("One real path: raw JSON -> warehouse star -> a chart a buyer reads", 26, WHITE, True)]])

    stages = [
        ("sources", "3 SOURCES", "Ticketmaster\nGoogle Trends\nYouTube"),
        ("bronze", "BRONZE", "GCS raw JSON\ndt-partitioned\nappend-only"),
        ("silver", "SILVER", "BQ fact constellation\n+ conformed dims"),
        ("gold", "GOLD", "fact_event_demand\n+ forecast_event_price"),
        ("app", "LIVE APP", "FastAPI + React\nprice + forecast\nlocal interest"),
    ]
    n = len(stages)
    gap = Inches(0.28)
    total_gap = gap * (n - 1)
    bw = (SW - 2 * MARGIN - total_gap) / n
    bh = Inches(2.5)
    top = Inches(2.4)
    x = MARGIN
    for key, head, body in stages:
        _box(slide, x, top, bw, Inches(0.55), fill=LAYER[key], shape=MSO_SHAPE.RECTANGLE)
        _text(slide, x, top + Inches(0.06), bw, Inches(0.45),
              [[(head, 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _box(slide, x, top + Inches(0.55), bw, bh - Inches(0.55), fill=LIGHT, line=LAYER[key],
             shape=MSO_SHAPE.RECTANGLE)
        _text(slide, x + Inches(0.05), top + Inches(0.7), bw - Inches(0.1), bh - Inches(0.7),
              [[(line, 13, INK, False)] for line in body.split("\n")],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if x + bw < SW - MARGIN - bw / 2:
            ar = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.02),
                                        top + bh / 2 - Inches(0.12), gap - Inches(0.04), Inches(0.24))
            ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT; ar.line.fill.background()
            ar.shadow.inherit = False
        x += bw + gap

    _box(slide, MARGIN, Inches(5.35), SW - 2 * MARGIN, Inches(1.4), fill=RGBColor(0xEC, 0xF3, 0xF7),
         line=ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    _text(slide, MARGIN + Inches(0.25), Inches(5.5), SW - 2 * MARGIN - Inches(0.5), Inches(1.1),
          [[("We follow ONE hero show the whole way: ", 16, INK, True),
            ("Everclear @ The Independent (SF).", 16, ACCENT, True)],
           [("Recognizable act, small ~500-cap room — the thesis in one show: a locally-known "
             "artist in a relatively small venue. 16 priced snapshots, real $136–236 band.", 14,
             MUTED, False)]])
    _notes(slide, """
~20s. This is the map you point back to at every step. Say: "data flows left to right; at each
step it gets cleaner and more joined, and it ends in something a non-technical user can read."
Name the hero show now (Everclear) so the next four slides have a single thread. Backup hero
show if Everclear's coverage shifts: The Wallflowers (nationwide). See hero-shows.md.
""")


def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH

    title_slide(prs)

    _content_slide(
        prs, "Quick recap · 30s · (assumed known)",
        "The question: will this show pop — early enough to act?",
        [("Domain: Bay-Area / nationwide live-music events", 0),
         ("Thesis: demand is high when a LOCALLY-popular artist plays a RELATIVELY SMALL venue", 0),
         ("Goal: help a buyer / seller / analyst anticipate sell-out + price pressure, early", 0),
         ("Assume you saw the midterm — here's only what's new", 1)],
        """
30s. Fast. Don't re-explain the whole domain. Land the thesis sentence clearly — every later
slide ties back to it. The "early enough to act" framing motivates the forecast + the app.
""")

    _content_slide(
        prs, "Schema — what changed · 30s · 5 pts",
        "Core schema unchanged — we added signal + a forecast table",
        [("Core schema is the SAME as the midterm: silver fact-constellation -> gold star "
          "(fact_event_demand). Say so plainly — the rubric rewards explaining 'no change' too", 0),
         ("Added since midterm:", 0),
         ("fact_trends_daily — daily-trajectory fact (1.58M rows) + 250-artist roster retarget", 1),
         ("forecast_event_price — model precomputed into gold (deterministic, no model-at-request)", 1),
         ("YouTube popularity signal (global artist reach)", 1),
         ("Why: deterministic / re-runnable gold; richer geo + time signal for the demand thesis", 0)],
        """
30s. The grader gives 5 pts for clearly explaining schema change OR no-change. Be explicit:
the STAR is unchanged; we enriched it. Mention the forecast is materialized into gold (a real
table), which sets up the gold + end-product slides.
""")

    pipeline_slide(prs)  # slide 4

    _content_slide(
        prs, "Bronze · ~50s · 15 pts",
        "Raw, untouched — exactly what the APIs return",
        [("Ticketmaster JSON — id, name, dates.start, _embedded.venues, "
          "priceRanges[].min/max/currency  (6 captures/day per state, append-only)", 0),
         ("= which shows exist, where, when, at what face price (the spine + the price we predict)", 1),
         ("Google Trends JSON — geoCode (DMA), value 0–100, query = artist, date", 0),
         ("= how big the artist is IN THIS METRO (geographic distribution of interest)", 1),
         ("YouTube JSON — subscriberCount, viewCount, videoCount per artist/day", 0),
         ("= how big the artist is OVERALL (global popularity magnitude)", 1),
         ("All land in GCS, dt-partitioned, append-only — a replayable audit log (~7.5 GB total)", 0)],
        """
~50s. Open an actual raw file live if you can (gsutil cat | head, or the GCS console) — showing
real JSON sells 'this is real data, not a fixture.' One business sentence per source; don't read
field lists. The point: three different lenses on demand (price, local interest, global reach).
""")

    _content_slide(
        prs, "Silver · ~70s · 15 pts",
        "Clean + conform into a fact constellation",
        [("fact_ticketmaster (event × snapshot_date) — parse priceRanges -> price_min/max; derive "
          "days_to_show; one row per (event, day); deterministic surrogate keys  [dbt model]", 0),
         ("fact_trends (artist × DMA × day) — interest 0–100, per-pull normalized", 0),
         ("Caveat aloud: comparable across time for one (artist, metro) — NEVER across artists/metros", 1),
         ("fact_youtube (artist × day) — subscribers + views", 0),
         ("Conformed dims = the join glue: dim_geo (Nielsen DMA — the only geography TM venues and "
          "Trends share; venue->DMA ~99.5%), dim_date, dim_artist, dim_venue, dim_event, bridge_event_artist", 0),
         ("Why a constellation: each fact keeps its NATIVE grain instead of flattening prematurely", 0)],
        """
~70s. This is where 'transformations defined by business logic' is earned. Emphasize: DMA is the
shared geography both sources resolve to — that's why the join is possible at all. The Trends
0–100 caveat is a credibility point: it shows you understand the data, not just plumbed it.
""")

    _content_slide(
        prs, "Transformation demo · ~50s · 15 pts (the required live transform)",
        "Run the gold join live — three feeds become one row",
        [("Run: dbt build -m fact_event_demand   (transforms execute IN the warehouse — ELT)", 0),
         ("Then bq query the hero show -> show TM price + joined local_interest + yt_subscribers "
          "appear in ONE wide row", 0),
         ("Say it while it runs: 'spine = every observed event-day; we LEFT-join local search "
          "interest + global YouTube reach onto it — no event-day is ever dropped'", 0),
         ("No-row-drop invariant is GX-tested (tests/assert_gold_rows_eq_spine.sql)", 1)],
        """
~50s. The rubric explicitly requires >=1 transformation DEMOED — this is that moment. Rehearse so
it returns in seconds (or have a pre-warmed result ready). If dbt is slow live, run the equivalent
BigQuery LEFT JOIN inline. Pick a hero show whose row is fully populated (see hero-shows.md).
""")

    _content_slide(
        prs, "Gold · ~60s · 15 pts",
        "The star + the forecast — what the app reads",
        [("fact_event_demand — one wide model-ready row per event per snapshot (609K rows / 9,657 events)", 0),
         ("Spine kept whole (no-row-drop, GX-tested); Trends + YouTube LEFT-joined on headliner + DMA", 1),
         ("forecast_event_price — demand/price model precomputed into gold (494,961 rows)", 0),
         ("Deterministic, seeded HistGradientBoostingRegressor; pooled cross-sectional", 1),
         ("Honest design given ~2 weeks of history — pool across shows, not per-show ARIMA", 1),
         ("So what: three messy feeds -> one table a model AND a UI can both consume", 0)],
        """
~60s. Show the real hero-show row in BigQuery. Be honest about the model: with ~2 weeks of
snapshots you can't do per-show time-series, so you pool cross-sectionally — that's a deliberate,
defensible choice, not a limitation you're hiding. Forecast is a TABLE (precomputed), so the app
never runs a model at request time.
""")

    _content_slide(
        prs, "Data quality · ~20s · 15 pts (optional)",
        "Great Expectations gates between every layer",
        [("15 checkpoints across bronze -> silver -> gold -> forecast", 0),
         ("Gold-refresh job FAILS FAST if the spine is dropped or row counts go wrong", 0),
         ("Quality is enforced in code, not eyeballed — strengthens 'transformations by business logic'", 0)],
        """
~20s. Optional — drop if you're tight on time. But it's a cheap credibility win and sets up the
tech-stack slide. One line: 'the pipeline refuses to publish bad data.'
""")

    _content_slide(
        prs, "END PRODUCT · ~90s · 10 pts · the consumable output",
        "One screen a non-technical buyer reads and acts on",
        [("Open the app -> pick a hero show from the dropdown (populated from fact_event_demand)", 0),
         ("Price history + forecast line (fact_event_demand + forecast_event_price)", 0),
         ("Local interest by metro / over time (fact_trends_daily)", 0),
         ("Artist YouTube reach (fact_youtube)", 0),
         ("Takeaway: 'Popular here? Small room? Price trending up as the show nears? -> act now or wait'", 0),
         ("Closes the loop back to the opening thesis", 1)],
        """
~90s — this is the 10-pt moment; spend real time here. The grader wants a consumable output a
non-technical person gets value from. Drive the live app on a frozen hero show. CRITICAL: have a
recorded screen capture / screenshots in assets/ as a fallback if the live demo fails. If the app
isn't ready, the Looker Studio fallback (on gold + forecast) still satisfies the rubric.
""")

    _content_slide(
        prs, "Tech stack · 1–2 min · 10 pts",
        "Each tool earns its place (expect a follow-up probe)",
        [("GCS — cheap durable raw landing  |  BigQuery — serverless warehouse for silver/gold", 0),
         ("dbt-bigquery (ELT) — transforms run IN the warehouse; one place for DDL, incremental "
          "materialization, tests, lineage.  Alt: in-memory pandas joins (don't scale)", 0),
         ("Cloud Run + Cloud Scheduler — scale-to-zero scheduled compute.  Alt: Airflow (deferred; bonus)", 0),
         ("Great Expectations — declarative DQ gates.  Alt: hand-rolled asserts", 0),
         ("scikit-learn HistGradientBoostingRegressor — deterministic; handles categoricals; fits pooled design", 0),
         ("FastAPI + React/Vite — live API over gold + UI  |  Terraform — reproducible GCP IaC", 0),
         ("GitHub Actions CI — ruff + pytest + terraform validate + GX smoke", 0)],
        """
1–2 min. The instructor WILL ask 'why this tool, what did you consider' for one component. The
most likely probe is dbt — be ready: ELT pushes joins/aggregations into BigQuery so it scales as
snapshots grow; one place for SQL, incremental builds, tests, lineage; vs. in-memory pandas that
doesn't scale. Have a one-line alternative-considered for every box on this slide.
""")

    _content_slide(
        prs, "Budget & cost · 1–2 min · 5 pts · CONFIRM WITH REAL BILLING",
        "~$5–20 / month at current volume (stated assumptions)",
        [("GCS storage ~7.5 GB (TM dominates) @ ~$0.02/GB-mo -> <$1/mo; grows with 6×/day TM captures", 0),
         ("BigQuery storage ~1–2 GB -> <$1/mo", 0),
         ("BigQuery query $6.25/TB on-demand; partition + cluster keep scans tiny -> single-digit $/mo", 0),
         ("Cloud Run + Scheduler — scale-to-zero -> ~free-tier to a few $/mo", 0),
         ("Secret Manager / monitoring -> negligible", 0),
         ("ROUGH TOTAL ~$5–20/mo (~$60–240/yr); dominated by storage growth + query frequency", 0),
         ("Every number is an assumption — replace with GCP billing-console actuals before the talk", 1)],
        """
1–2 min. The rubric wants realistic estimates WITH stated assumptions — so say the assumptions out
loud (data volume, query frequency, scale-to-zero). REPLACE these order-of-magnitude numbers with
real billing-console figures before presenting. Tie cost to the drivers: storage grows with the
6×/day TM sweep; query cost grows with app/query frequency.
""")

    _content_slide(
        prs, "Close · bonus",
        "Built, working — and we shipped orchestration",
        [("Recap in one line: three live feeds -> medallion pipeline -> one model-ready table -> a "
          "consumable app", 0),
         ("Next: wire the daily Trends trajectory into gold; deploy hardening", 0),
         ("BONUS: we built real orchestration — a Cloud Run gold-refresh job chaining "
          "silver -> dbt build -> forecast -> GX gate, daily", 0),
         ("Per the assignment, orchestration already built can count toward bonus — let's talk", 1)],
        """
Close strong + claim the bonus. Orchestration was made optional, but you built it (Cloud Run
gold-refresh). The assignment says if you already built orchestration, talk to the instructor for
bonus points — so explicitly flag it here. End on the thesis -> product loop being closed.
""")

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "event-demand-final-deck.pptx")
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    build()
